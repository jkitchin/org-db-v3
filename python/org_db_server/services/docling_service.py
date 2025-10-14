"""Service for converting linked files to markdown using docling."""
import logging
import hashlib
import os
import sys
import subprocess
import json
from pathlib import Path
from typing import Optional, Dict, Any

logger = logging.getLogger(__name__)

# Enable HuggingFace downloads for PDF processing models
# PDF processing requires downloading AI models from HuggingFace Hub
if 'HF_HUB_OFFLINE' not in os.environ:
    os.environ['HF_HUB_OFFLINE'] = '0'


class DoclingService:
    """Service for converting documents to markdown using docling."""

    # Supported file extensions
    SUPPORTED_EXTENSIONS = {
        # PDFs (handled by pymupdf4llm - lightweight and fast)
        '.pdf',
        # DOCX (handled by python-docx - lightweight and fast)
        '.docx',
        # PPTX (handled by python-pptx - lightweight and fast)
        '.pptx',
        # Office documents (handled by docling in subprocess - last resort)
        '.doc', '.xlsx', '.xls', '.ppt',
        # Web/markup formats
        '.html', '.htm', '.xhtml', '.md', '.markdown', '.asciidoc', '.adoc',
        # Data formats
        '.csv',
        # Image formats (with OCR)
        '.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.webp',
        # Specialized formats
        '.vtt', '.xml', '.json',
    }

    def __init__(self, use_subprocess: bool = True):
        """Initialize the docling service.

        Args:
            use_subprocess: If True, run Docling in isolated subprocess (safer).
                           If False, run in-process (faster but can crash server).
        """
        self.use_subprocess = use_subprocess
        self._worker_path = Path(__file__).parent / "docling_worker.py"
        logger.info(f"DoclingService initialized (subprocess={use_subprocess})")

    @staticmethod
    def calculate_md5(file_path: str) -> str:
        """Calculate MD5 hash of a file for change detection."""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def is_supported(self, file_path: str) -> bool:
        """Check if a file extension is supported by docling."""
        ext = Path(file_path).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS

    def convert_to_markdown(
        self,
        file_path: str,
        max_file_size: int = 52428800  # 50MB default
    ) -> Dict[str, Any]:
        """
        Convert a document to markdown using docling.

        Args:
            file_path: Path to the file to convert
            max_file_size: Maximum file size in bytes (default 50MB)

        Returns:
            Dictionary with:
                - status: 'success', 'error', 'skipped'
                - markdown: The converted markdown text (if successful)
                - md5: MD5 hash of the source file
                - error: Error message (if failed)
                - file_size: Size of the file in bytes
        """
        path = Path(file_path)

        # Check file exists
        if not path.exists():
            logger.warning(f"File not found: {file_path}")
            return {
                "status": "error",
                "error": "File not found",
                "md5": None,
                "file_size": 0
            }

        # Check file size
        file_size = path.stat().st_size
        if file_size > max_file_size:
            logger.warning(f"File too large: {file_path} ({file_size} bytes)")
            return {
                "status": "skipped",
                "error": f"File too large: {file_size / 1024 / 1024:.1f}MB",
                "md5": None,
                "file_size": file_size
            }

        # Check extension
        if not self.is_supported(file_path):
            ext = path.suffix
            logger.warning(f"Unsupported file extension: {ext}")
            return {
                "status": "error",
                "error": f"Unsupported file extension: {ext}",
                "md5": None,
                "file_size": file_size
            }

        # Calculate MD5
        try:
            md5 = self.calculate_md5(file_path)
        except Exception as e:
            logger.error(f"Error calculating MD5 for {file_path}: {e}")
            return {
                "status": "error",
                "error": f"Error calculating MD5: {str(e)}",
                "md5": None,
                "file_size": file_size
            }

        # Convert document
        ext = path.suffix.lower()

        # Use lightweight libraries for common formats (fast, low memory)
        if ext == '.pdf':
            return self._convert_pdf_with_pymupdf(file_path, md5, file_size)
        elif ext == '.docx':
            return self._convert_docx_with_python_docx(file_path, md5, file_size)
        elif ext == '.pptx':
            return self._convert_pptx_with_python_pptx(file_path, md5, file_size)
        elif self.use_subprocess:
            # Run in isolated subprocess for safety (docling for other formats)
            return self._convert_in_subprocess(file_path, max_file_size, md5, file_size)
        else:
            # Run in-process (legacy, not recommended)
            return self._convert_in_process(file_path, md5, file_size)

    def _convert_pdf_with_pymupdf(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert PDF using lightweight pymupdf4llm (fast, low memory)."""
        try:
            import pymupdf4llm

            logger.info(f"Converting PDF {file_path} with pymupdf4llm...")
            markdown_text = pymupdf4llm.to_markdown(file_path)

            logger.info(f"Successfully converted {file_path} ({len(markdown_text)} chars)")
            return {
                "status": "success",
                "markdown": markdown_text,
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting PDF {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_docx_with_python_docx(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert DOCX using lightweight python-docx (fast, low memory)."""
        try:
            from docx import Document

            logger.info(f"Converting DOCX {file_path} with python-docx...")
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

            logger.info(f"Successfully converted {file_path} ({len(text)} chars)")
            return {
                "status": "success",
                "markdown": text,  # Plain text, not markdown, but good enough
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting DOCX {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_pptx_with_python_pptx(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert PPTX using lightweight python-pptx (fast, low memory)."""
        try:
            from pptx import Presentation

            logger.info(f"Converting PPTX {file_path} with python-pptx...")
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

            result_text = "\n".join(text)

            logger.info(f"Successfully converted {file_path} ({len(result_text)} chars)")
            return {
                "status": "success",
                "markdown": result_text,  # Plain text, not markdown, but good enough
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting PPTX {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_in_subprocess(
        self,
        file_path: str,
        max_file_size: int,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert file using isolated subprocess (safe from crashes)."""
        try:
            logger.info(f"Converting {file_path} in subprocess...")

            # Run worker process with timeout
            cmd = [sys.executable, str(self._worker_path), file_path, str(max_file_size)]

            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=300,  # 5 minute timeout
                env={**os.environ, 'HF_HUB_OFFLINE': '0'},
                start_new_session=True  # Isolate from parent process group
            )

            # Parse JSON output
            try:
                output = json.loads(result.stdout)
            except json.JSONDecodeError as e:
                logger.error(f"Failed to parse worker output: {result.stdout}")
                logger.error(f"Worker stderr: {result.stderr}")
                return {
                    "status": "error",
                    "error": f"Worker output parse error: {str(e)}",
                    "md5": md5,
                    "file_size": file_size
                }

            # Check if worker crashed
            if result.returncode not in (0, 1):
                logger.error(f"Worker crashed with code {result.returncode}")
                logger.error(f"Worker stderr: {result.stderr}")
                return {
                    "status": "error",
                    "error": f"Worker crashed (code {result.returncode}): {output.get('error', 'unknown error')}",
                    "md5": md5,
                    "file_size": file_size
                }

            # Worker succeeded or failed gracefully
            if output["status"] == "success":
                logger.info(f"Successfully converted {file_path} ({len(output['markdown'])} chars)")

            return output

        except subprocess.TimeoutExpired:
            logger.error(f"Conversion timeout for {file_path}")
            return {
                "status": "error",
                "error": "Conversion timeout (5 minutes)",
                "md5": md5,
                "file_size": file_size
            }
        except Exception as e:
            logger.error(f"Error in subprocess conversion: {e}", exc_info=True)
            return {
                "status": "error",
                "error": f"Subprocess error: {str(e)}",
                "md5": md5,
                "file_size": file_size
            }

    def _convert_in_process(self, file_path: str, md5: str, file_size: int) -> Dict[str, Any]:
        """Convert file in-process (legacy, can crash server)."""
        try:
            from docling.document_converter import DocumentConverter

            logger.info(f"Converting {file_path} in-process (UNSAFE)...")
            converter = DocumentConverter()
            result = converter.convert(file_path)

            # Export to markdown
            markdown_text = result.document.export_to_markdown()

            # Force cleanup
            del result
            del converter
            import gc
            gc.collect()

            logger.info(f"Successfully converted {file_path} ({len(markdown_text)} chars)")
            return {
                "status": "success",
                "markdown": markdown_text,
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}", exc_info=True)
            import gc
            gc.collect()
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }


# Global instance
_docling_service: Optional[DoclingService] = None


def get_docling_service() -> DoclingService:
    """Get or create the global docling service instance."""
    global _docling_service
    if _docling_service is None:
        _docling_service = DoclingService()
    return _docling_service
